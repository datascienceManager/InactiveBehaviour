"""
CHURN ANALYSIS POWERPOINT GENERATOR
====================================

This script creates a comprehensive PowerPoint presentation from the
R-generated churn analysis outputs (tables and visualizations).

Prerequisites:
1. Run master_churn_workflow.R to generate all analysis outputs
2. Ensure churn_visuals/ and churn_tables/ directories exist
3. Install: pip install python-pptx

Author: Data Analytics Team
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pathlib import Path
import os


class ChurnPresentationBuilder:
    """Build comprehensive churn analysis PowerPoint presentation"""
    
    def __init__(self, visuals_dir="churn_visuals", tables_dir="churn_tables"):
        """Initialize the presentation builder"""
        self.visuals_dir = Path(visuals_dir)
        self.tables_dir = Path(tables_dir)
        self.prs = Presentation()
        self.prs.slide_width = Inches(10)
        self.prs.slide_height = Inches(7.5)
        
        # Brand colors
        self.brand_blue = RGBColor(68, 114, 196)
        self.brand_orange = RGBColor(237, 125, 49)
        
        # Verify directories
        self._check_directories()
    
    def _check_directories(self):
        """Check if required directories exist"""
        if not self.visuals_dir.exists():
            print(f"⚠️  Warning: {self.visuals_dir} not found")
        if not self.tables_dir.exists():
            print(f"⚠️  Warning: {self.tables_dir} not found")
    
    def add_title_slide(self, title, subtitle=""):
        """Add title slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # Title
        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(2.5), Inches(8), Inches(1.5)
        )
        title_frame = title_box.text_frame
        title_frame.text = title
        p = title_frame.paragraphs[0]
        p.font.size = Pt(48)
        p.font.bold = True
        p.font.color.rgb = self.brand_blue
        p.alignment = PP_ALIGN.CENTER
        
        # Subtitle
        if subtitle:
            subtitle_box = slide.shapes.add_textbox(
                Inches(1), Inches(4.2), Inches(8), Inches(1)
            )
            subtitle_frame = subtitle_box.text_frame
            subtitle_frame.text = subtitle
            p = subtitle_frame.paragraphs[0]
            p.font.size = Pt(24)
            p.font.color.rgb = self.brand_orange
            p.alignment = PP_ALIGN.CENTER
        
        return slide
    
    def add_section_divider(self, section_title):
        """Add section divider slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # Section title
        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(3), Inches(8), Inches(1.5)
        )
        title_frame = title_box.text_frame
        title_frame.text = section_title
        p = title_frame.paragraphs[0]
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = self.brand_blue
        p.alignment = PP_ALIGN.CENTER
        
        return slide
    
    def add_image_slide(self, title, image_name, directory="visuals",
                       left=0.5, top=1.5, width=9):
        """Add slide with title and image"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.4), Inches(9), Inches(0.8)
        )
        title_frame = title_box.text_frame
        title_frame.text = title
        p = title_frame.paragraphs[0]
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = self.brand_blue
        
        # Image
        img_dir = self.visuals_dir if directory == "visuals" else self.tables_dir
        image_path = img_dir / image_name
        
        if image_path.exists():
            slide.shapes.add_picture(
                str(image_path),
                Inches(left),
                Inches(top),
                width=Inches(width)
            )
            print(f"  ✓ Added: {image_name}")
        else:
            print(f"  ⚠️  Missing: {image_name}")
            # Add placeholder text
            text_box = slide.shapes.add_textbox(
                Inches(2), Inches(3), Inches(6), Inches(1)
            )
            text_frame = text_box.text_frame
            text_frame.text = f"Image not found: {image_name}"
            text_frame.paragraphs[0].font.size = Pt(16)
            text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 0, 0)
        
        return slide
    
    def add_split_slide(self, title, left_image, right_image,
                       left_dir="visuals", right_dir="visuals"):
        """Add slide with two images side by side"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.4), Inches(9), Inches(0.8)
        )
        title_frame = title_box.text_frame
        title_frame.text = title
        p = title_frame.paragraphs[0]
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = self.brand_blue
        
        # Left image
        left_path = (self.visuals_dir if left_dir == "visuals" else self.tables_dir) / left_image
        if left_path.exists():
            slide.shapes.add_picture(
                str(left_path),
                Inches(0.4),
                Inches(1.8),
                width=Inches(4.5)
            )
        
        # Right image
        right_path = (self.visuals_dir if right_dir == "visuals" else self.tables_dir) / right_image
        if right_path.exists():
            slide.shapes.add_picture(
                str(right_path),
                Inches(5.2),
                Inches(1.8),
                width=Inches(4.5)
            )
        
        return slide
    
    def add_bullet_slide(self, title, bullets):
        """Add slide with bullet points"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.4), Inches(9), Inches(0.8)
        )
        title_frame = title_box.text_frame
        title_frame.text = title
        p = title_frame.paragraphs[0]
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = self.brand_blue
        
        # Bullets
        text_box = slide.shapes.add_textbox(
            Inches(1), Inches(1.8), Inches(8), Inches(5)
        )
        text_frame = text_box.text_frame
        text_frame.word_wrap = True
        
        for i, bullet in enumerate(bullets):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            
            p.text = bullet
            p.level = 0
            p.font.size = Pt(18)
            p.space_before = Pt(12)
        
        return slide
    
    def save(self, filename="Churn_Analysis_Report.pptx"):
        """Save the presentation"""
        self.prs.save(filename)
        print(f"\n✓ Presentation saved: {filename}")
        print(f"✓ Total slides: {len(self.prs.slides)}")


def create_comprehensive_presentation():
    """Create complete churn analysis presentation"""
    
    print("=" * 70)
    print("Creating Comprehensive Churn Analysis Presentation")
    print("=" * 70)
    print()
    
    builder = ChurnPresentationBuilder()
    
    # ====================================================================
    # SLIDE 1: Title Slide
    # ====================================================================
    print("Adding slides...")
    builder.add_title_slide(
        "OTT Platform Churn Analysis",
        "Predictive Modeling & Strategic Insights"
    )
    print("  ✓ Title slide")
    
    # ====================================================================
    # SLIDE 2: Agenda
    # ====================================================================
    builder.add_bullet_slide(
        "Agenda",
        [
            "Executive Summary & Key Metrics",
            "Exploratory Data Analysis",
            "Churn Drivers & Patterns",
            "Predictive Model Results",
            "Risk Segmentation",
            "Recommendations & Action Plan"
        ]
    )
    print("  ✓ Agenda slide")
    
    # ====================================================================
    # SECTION 1: Executive Summary
    # ====================================================================
    builder.add_section_divider("Executive Summary")
    print("  ✓ Section divider: Executive Summary")
    
    builder.add_image_slide(
        "Key Metrics Overview",
        "01_executive_summary.png",
        directory="tables",
        left=1.5,
        width=7
    )
    
    builder.add_split_slide(
        "Churn Overview",
        "01_overall_churn_rate.png",
        "05_churn_type_breakdown.png"
    )
    
    # ====================================================================
    # SECTION 2: Exploratory Analysis
    # ====================================================================
    builder.add_section_divider("Exploratory Data Analysis")
    print("  ✓ Section divider: EDA")
    
    builder.add_image_slide(
        "Churn by Product Type",
        "02_churn_by_product.png"
    )
    
    builder.add_image_slide(
        "Churn by Customer Tenure",
        "03_churn_by_tenure.png"
    )
    
    builder.add_image_slide(
        "Churn by Acquisition Channel",
        "06_churn_by_channel.png"
    )
    
    builder.add_image_slide(
        "Product Performance Analysis",
        "02_product_analysis.png",
        directory="tables",
        left=0.8,
        width=8.5
    )
    
    # ====================================================================
    # SECTION 3: Churn Drivers
    # ====================================================================
    builder.add_section_divider("Key Churn Drivers")
    print("  ✓ Section divider: Churn Drivers")
    
    builder.add_image_slide(
        "Engagement Impact on Churn",
        "04_engagement_distribution.png"
    )
    
    builder.add_image_slide(
        "Content Viewing Preferences",
        "09_content_preferences.png"
    )
    
    builder.add_image_slide(
        "Cohort Retention Trends",
        "07_cohort_retention.png"
    )
    
    builder.add_image_slide(
        "RFM Customer Segmentation",
        "08_rfm_segmentation.png"
    )
    
    # ====================================================================
    # SECTION 4: Predictive Model
    # ====================================================================
    builder.add_section_divider("Predictive Modeling Results")
    print("  ✓ Section divider: Predictive Modeling")
    
    builder.add_image_slide(
        "Model Performance Comparison",
        "03_model_comparison.png",
        directory="tables",
        left=1.5,
        width=7
    )
    
    builder.add_image_slide(
        "ROC Curves - Model Comparison",
        "11_roc_curves.png",
        left=0.5
    )
    
    builder.add_image_slide(
        "Top 15 Most Important Features",
        "12_feature_importance.png",
        left=0.5
    )
    
    builder.add_image_slide(
        "Feature Importance - Detailed",
        "05_feature_importance.png",
        directory="tables",
        left=1.5,
        width=7
    )
    
    builder.add_image_slide(
        "Model Calibration",
        "13_calibration_plot.png",
        left=0.5
    )
    
    # ====================================================================
    # SECTION 5: Risk Segmentation
    # ====================================================================
    builder.add_section_divider("Risk Segmentation & Targeting")
    print("  ✓ Section divider: Risk Segmentation")
    
    builder.add_image_slide(
        "Churn Risk Score Distribution",
        "10_risk_score_distribution.png"
    )
    
    builder.add_image_slide(
        "Risk Segment Validation",
        "14_risk_validation.png"
    )
    
    builder.add_image_slide(
        "Risk Segmentation Analysis",
        "04_risk_segmentation.png",
        directory="tables",
        left=0.8,
        width=8.5
    )
    
    # ====================================================================
    # SECTION 6: Recommendations
    # ====================================================================
    builder.add_section_divider("Strategic Recommendations")
    print("  ✓ Section divider: Recommendations")
    
    builder.add_bullet_slide(
        "Immediate Actions - High Priority",
        [
            "Target very high-risk customers with personalized retention offers",
            "Launch re-engagement campaign for inactive users (>30 days no viewing)",
            "Prioritize high-value customers showing declining engagement",
            "Address product-specific issues (highest churn products)",
            "Implement early warning system using the predictive model"
        ]
    )
    
    builder.add_bullet_slide(
        "Product & Content Improvements",
        [
            "Review and enhance high-churn product offerings",
            "Improve content recommendations based on viewing patterns",
            "Develop targeted content for at-risk customer segments",
            "Optimize user experience on high-churn channels",
            "Create exclusive content for loyal customers"
        ]
    )
    
    builder.add_bullet_slide(
        "Engagement Strategies",
        [
            "Personalized content suggestions based on viewing history",
            "Implement viewing milestones and rewards program",
            "Send proactive retention communications to at-risk customers",
            "Create winback campaigns for churned customers",
            "Develop onboarding improvements for new subscribers"
        ]
    )
    
    builder.add_bullet_slide(
        "Monitoring & Optimization",
        [
            "Weekly risk scoring of all active subscribers",
            "Monitor churn rates by segment and take corrective actions",
            "A/B test retention interventions and measure impact",
            "Track model performance and retrain quarterly",
            "Establish early warning KPIs and dashboards"
        ]
    )
    
    # ====================================================================
    # SECTION 7: Business Impact
    # ====================================================================
    builder.add_section_divider("Expected Business Impact")
    print("  ✓ Section divider: Business Impact")
    
    builder.add_bullet_slide(
        "Projected Outcomes",
        [
            "Reduce overall churn rate by 15-20% through targeted interventions",
            "Save $XXX,XXX in annual recurring revenue",
            "Improve customer lifetime value by extending average tenure",
            "Increase retention rate for high-value customers by 25%",
            "Achieve 85%+ model accuracy in identifying at-risk customers"
        ]
    )
    
    # ====================================================================
    # SECTION 8: Next Steps
    # ====================================================================
    builder.add_bullet_slide(
        "Implementation Roadmap - Next 90 Days",
        [
            "Week 1-2: Deploy model in production, score all active customers",
            "Week 3-4: Design and launch retention campaigns for high-risk segments",
            "Week 5-6: Implement product improvements based on churn analysis",
            "Week 7-8: Launch content recommendation enhancements",
            "Week 9-12: Monitor results, A/B test interventions, optimize approach"
        ]
    )
    
    # ====================================================================
    # Final Slide
    # ====================================================================
    builder.add_title_slide(
        "Questions?",
        "Thank you for your attention"
    )
    print("  ✓ Closing slide")
    
    # Save presentation
    builder.save("Churn_Analysis_Report.pptx")
    
    print()
    print("=" * 70)
    print("Presentation created successfully!")
    print("=" * 70)


if __name__ == "__main__":
    # Check if required directories exist
    if not Path("churn_visuals").exists() and not Path("churn_tables").exists():
        print("⚠️  ERROR: Required directories not found!")
        print()
        print("Please run the R analysis first:")
        print("  Rscript master_churn_workflow.R")
        print()
        print("This will generate:")
        print("  - churn_visuals/ directory with all charts")
        print("  - churn_tables/ directory with all tables")
        print()
    else:
        create_comprehensive_presentation()
        
        print()
        print("NEXT STEPS:")
        print("=" * 70)
        print("1. Open 'Churn_Analysis_Report.pptx'")
        print("2. Review and customize slides as needed")
        print("3. Add company branding and specific metrics")
        print("4. Present to stakeholders")
        print()
        print("=" * 70)
